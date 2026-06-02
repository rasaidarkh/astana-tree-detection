# -*- coding: utf-8 -*-
"""Generate an editable AITU-style, system-first defense deck as .pptx.

Mirrors deck_aitu_system.html but as native PowerPoint shapes/text/tables so it
imports into Canva fully editable. Run with the repo venv:
    venv/Scripts/python.exe thesis/slides/build_pptx.py
"""
import os, sys
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
FIG  = os.path.join(os.path.dirname(HERE), "figures")
ASSET= os.path.join(HERE, "assets")

# ---- palette ----
BLUE   = RGBColor(0x2E,0x6D,0xB4)
CYAN   = RGBColor(0x29,0xAB,0xE2)
NAVY   = RGBColor(0x17,0x3E,0x63)
INK    = RGBColor(0x1A,0x1A,0x1A)
INK2   = RGBColor(0x33,0x41,0x4D)
MUTED  = RGBColor(0x6A,0x76,0x82)
CARD   = RGBColor(0xEE,0xF4,0xFB)
CARDLN = RGBColor(0xCF,0xE0,0xF1)
LINE   = RGBColor(0xD7,0xDE,0xE4)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
BUILT  = RGBColor(0xDC,0xEB,0xFA)
GREY   = RGBColor(0xF2,0xF4,0xF6)
SERIF  = "Times New Roman"

EMU_IN = 914400
SW, SH = 13.333, 7.5
ML, MR = 0.58, 0.58
CW = SW - ML - MR  # content width

RATIO = {
    "yolo_all_runs_map50.png":2.18, "model_comparison_barchart.png":2.34,
    "cross_model_8way_botanical.jpg":2.05, "ui_canopy_map_view.png":1.60,
    "ui_canopy_scan_model_modal.png":1.60, "ui_canopy_library_modal.png":1.60,
    "ui_canopy_map_heat.png":1.60, "yolo_v2_finetune_val_4tile_strip.png":3.35,
    "ui_canopy_botanical_park.jpg":1.76,
}

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------
def slide():
    return prs.slides.add_slide(BLANK)

def _noline(shp):
    shp.line.fill.background()

def rect(s, x, y, w, h, fill, line=None, line_w=0.75):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.shadow.inherit = False
    if line is None:
        _noline(sp)
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp

def band(s):
    # gradient ribbon navy -> cyan; fallback to solid navy
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.6))
    b.shadow.inherit = False; _noline(b)
    try:
        b.fill.gradient()
        stops = b.fill.gradient_stops
        stops[0].color.rgb = NAVY;  stops[0].position = 0.0
        stops[1].color.rgb = CYAN;  stops[1].position = 1.0
        b.fill.gradient_angle = 0.0
    except Exception:
        b.fill.solid(); b.fill.fore_color.rgb = NAVY
    # two-tone lighter block on the right
    r = rect(s, SW-1.55, 0, 1.55, 0.6, CYAN);
    # shield
    shx, shw, shy = 0.62, 0.8, 0.0
    topH, ptH = 0.78, 0.30
    rect(s, shx, shy, shw, topH, CYAN)
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(shx), Inches(shy+topH), Inches(shw), Inches(ptH))
    tri.rotation = 180; tri.fill.solid(); tri.fill.fore_color.rgb = CYAN; tri.shadow.inherit=False; _noline(tri)
    # logo: white ring + dot
    ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(shx+shw/2-0.2), Inches(0.12), Inches(0.4), Inches(0.4))
    ring.fill.background(); ring.line.color.rgb = WHITE; ring.line.width = Pt(2); ring.shadow.inherit=False
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(shx+shw/2-0.06), Inches(0.30), Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = WHITE; _noline(dot); dot.shadow.inherit=False

def footer(s, left, page):
    ln = rect(s, ML, SH-0.62, CW, 0.012, LINE)
    tb(s, ML, SH-0.55, CW-0.6, 0.3, [(left, 9, MUTED, False, True)], anchor=MSO_ANCHOR.TOP)
    tb(s, SW-MR-0.6, SH-0.55, 0.6, 0.3, [(str(page), 9, MUTED, False, True)], align=PP_ALIGN.RIGHT)

def tb(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, space=3, leading=None):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,(txt,sz,col,bold,ital) in enumerate(paras):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        if leading: p.line_spacing = leading
        r = p.add_run(); r.text = txt
        f = r.font; f.size=Pt(sz); f.bold=bold; f.italic=ital; f.name=SERIF; f.color.rgb=col
    return box

def rich(s, x, y, w, h, segments, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, leading=1.05):
    """segments: list of paragraphs; each paragraph is list of (text,size,col,bold,ital)."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,runs in enumerate(segments):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=leading; p.space_after=Pt(7); p.space_before=Pt(0)
        for (txt,sz,col,bold,ital) in runs:
            r=p.add_run(); r.text=txt; f=r.font
            f.size=Pt(sz); f.bold=bold; f.italic=ital; f.name=SERIF; f.color.rgb=col
    return box

def title(s, text, y=0.95, size=30, color=INK):
    tb(s, ML, y, CW, 0.8, [(text, size, color, True, False)])

def eyebrow(s, text, y=0.86):
    tb(s, ML, y, CW, 0.3, [(text, 13, BLUE, True, False)])

def card(s, x, y, w, h, ttl, body_segs, accent=BLUE, fill=CARD, tsize=14, bsize=12.5):
    rect(s, x, y, w, h, fill, line=CARDLN, line_w=0.75)
    rect(s, x, y, 0.06, h, accent)
    segs = []
    if ttl: segs.append([(ttl, tsize, BLUE, True, False)])
    for b in body_segs: segs.append(b)
    rich(s, x+0.18, y+0.12, w-0.34, h-0.24, segs, leading=1.08)

def place_img(s, fname, x, y, w, border=True):
    path = os.path.join(FIG, fname)
    r = RATIO.get(fname, 1.6)
    h = w / r
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    if border:
        pic.line.color.rgb = LINE; pic.line.width = Pt(1)
    return pic, w, h

def block(s, x, y, w, h, text, fill, outline=False):
    sp = rect(s, x, y, w, h, fill)
    if outline:
        sp.line.color.rgb = NAVY; sp.line.width = Pt(1.5)
    tf = sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_top=Pt(2); tf.margin_bottom=Pt(2)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; f=r.font; f.size=Pt(13); f.bold=True; f.name=SERIF; f.color.rgb=WHITE

# ---------------- slides ----------------
# 01 TITLE
s = slide(); band(s)
tb(s, 1.4, 1.7, SW-2.8, 2.6,
   [("Development of a Deep-Learning System for Automated Tree Recognition & Green-Space Mapping in Urban Environments",
     30, BLUE, True, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, leading=1.1)
tb(s, 1.0, 4.7, SW-2.0, 0.5, [("Rasul Aidarkhanov    ·    Berik Sharipov    ·    Anuar Totin", 19, INK, False, False)], align=PP_ALIGN.CENTER)
tb(s, 1.0, 5.18, SW-2.0, 0.4, [("Supervisor: Syndar Satbayev", 16, INK2, False, False)], align=PP_ALIGN.CENTER)
tb(s, 1.0, 6.0, SW-2.0, 0.4, [("Astana IT University   ·   6B06101 Computer Science   ·   2026", 14, MUTED, False, False)], align=PP_ALIGN.CENTER)

# 02 CONTENTS
s = slide(); band(s); title(s, "Contents")
toc = [("01","Relevance of the topic"),("02","The gap & research aim"),
       ("03","What this project is — a system"),("04","Objectives & methods"),
       ("05","Problem statement · Input / Output"),("06","Literature review & the gap"),
       ("07","The automated pipeline"),
       ("08","Data, database & hardware"),("09","How we measure success"),
       ("10","Four detection engines"),("11","Why several models · ensembles"),
       ("12","Results on one common set"),("13","Canopy — the system in use"),
       ("14","Discussion & conclusion")]
colL = toc[:7]; colR = toc[7:]
def toc_col(x, items):
    segs=[]
    for n,t in items:
        segs.append([(n+"   ", 17, BLUE, True, False),(t, 17, INK, False, False)])
    rich(s, x, 1.75, 5.6, 4.6, segs, leading=1.0)
    # bump space_after
toc_col(ML, colL); toc_col(ML+6.2, colR)
footer(s, "Presenter — Rasul Aidarkhanov · system + YOLO + web app", 2)

# 03 RELEVANCE
s = slide(); band(s); title(s, "Relevance", y=1.0)
rich(s, ML, 1.95, CW, 4.2, [
 [("Urban trees regulate microclimate, sequester carbon and mitigate the heat-island effect — so a growing city like Astana needs an ", 17, INK2, False, False),
  ("up-to-date, spatially accurate inventory", 17, INK, True, False),
  (" of its green infrastructure for planning and climate-adaptation decisions.", 17, INK2, False, False)],
 [("Traditional inventories rely on ", 17, INK2, False, False),("manual field surveys", 17, INK, True, False),
  (" — slow, expensive, and outdated before completion. A fundamental scalability problem for any large urban area.", 17, INK2, False, False)],
 [("Deep-learning detection on free satellite imagery is a feasible alternative — but the real gap is that ", 17, INK2, False, False),
  ("no automated process for this exists in Astana at all.", 17, BLUE, True, False),
  (" So the task we set ourselves was an engineering one: build that process.", 17, INK2, False, False)],
], leading=1.3)
footer(s, "Criterion 1 · Topic disclosure", 3)

# 04 GAP & AIM
s = slide(); band(s); eyebrow(s, "Why we couldn't just use an existing tool", y=1.5)
tb(s, ML, 1.85, 4.2, 1.6, [("0.012", 96, BLUE, True, False)])
rich(s, ML+4.4, 1.95, CW-4.4, 1.7, [
 [("Box mAP@50", 16, INK, True, False)],
 [("A popular ready-made detector (NEON DeepForest), run as-is on Astana. Essentially blind. There was no off-the-shelf system that worked here, so building one was the only option.", 15, INK2, False, False)],
], leading=1.2, anchor=MSO_ANCHOR.MIDDLE)
card(s, ML, 4.2, CW, 1.5, "Research aim",
     [[("Design, build and evaluate an ", 15, INK2, False, False),("automated end-to-end system", 15, INK, True, False),
       (" that recognises trees and maps green space from Astana satellite imagery — delivered as a working tool, with detection accuracy reported honestly.", 15, INK2, False, False)]],
     tsize=15, bsize=15)
footer(s, "Criterion 1 · Aim", 4)

# 05 SYSTEM NOT MODEL
s = slide(); band(s); eyebrow(s, "The thesis, in one idea")
title(s, "The deliverable is a system — not a model", y=1.18, size=30)
rich(s, ML, 2.0, CW, 0.9, [
 [("Our title says ", 15.5, INK2, False, False),('“automated … system.”', 15.5, INK, True, False),
  (" So the unit of work isn't “which neural net scores highest” — it's the whole machine that turns a raw satellite capture into a usable, georeferenced tree inventory. Everything else is a ", 15.5, INK2, False, False),
  ("component of building that:", 15.5, INK, True, False)]], leading=1.25)
comp = [("Dataset","What the engine learns from — a means, not the goal."),
        ("The models","Four interchangeable engines that plug into the pipeline."),
        ("The experiments","How we tuned the engine — evidence of depth, not the product."),
        ("UI & features","What makes it usable by a non-technical operator.")]
cw4 = (CW-3*0.2)/4
for i,(t,b) in enumerate(comp):
    card(s, ML+i*(cw4+0.2), 3.15, cw4, 1.35, t, [[(b,12,INK2,False,False)]], tsize=14, bsize=12)
card(s, ML, 4.75, CW, 1.25, "",
     [[("We could have shipped the bare minimum — one model, a mask, no metrics. Instead we built a ", 14, INK2, False, False),
       ("tool someone at a city service could actually operate.", 14, INK, True, False),
       (" The slides that follow lead with the system; the models come later, as parts of it.", 14, INK2, False, False)]], bsize=14)
footer(s, "Criterion 1 · Framing — engineering deliverable", 5)

# 06 OBJECTIVES & METHODS
s = slide(); band(s); title(s, "Objectives & methods")
objL = [("1.","Build an annotated Astana dataset (tiled satellite imagery, instance polygons)."),
        ("2.","Train and compare four detection engines — YOLOv8-seg, Mask R-CNN, DeepForest, SAM 2."),
        ("3.","Combine them through ensembles and test whether that helps.")]
objR = [("4.","Evaluate everything on one common test set for a fair comparison."),
        ("5.","Wrap it all in an automated web application — the deliverable that ties tasks 1–4 together.")]
def obj_col(x, items):
    segs=[[(n+"  ",16,BLUE,True,False),(t,16,INK2,False,False)] for n,t in items]
    rich(s, x, 1.9, 5.7, 2.6, segs, leading=1.15)
obj_col(ML, objL); obj_col(ML+6.2, objR)
card(s, ML, 4.7, CW, 1.2, "Methods",
     [[("Supervised deep learning · transfer learning from pretrained weights · sliding-window tiling · pixel → WGS-84 geo-referencing · mAP@50 as the single comparison metric · software engineering (FastAPI · React · SQLite).", 14, INK2, False, False)]], bsize=14)
footer(s, "Criterion 1 · Tasks & methods", 6)

# 07 PROBLEM STATEMENT
s = slide(); band(s); title(s, "Problem statement — Input / Output")
half=(CW-0.3)/2
card(s, ML, 1.85, half, 2.6, "Input",
     [[("A satellite capture of an area of Astana at zoom 19 (~0.3 m/px) — uploaded, or drawn on a map and auto-fetched from ", 13.5, INK2, False, False),
       ("ESRI or Google", 13.5, INK, True, False),(" tiles — cut into overlapping ", 13.5, INK2, False, False),
       ("640 px RGB tiles.", 13.5, INK, True, False)],
      [("Plus the operator's choices: model, confidence, geo-mode.", 13.5, INK2, False, False)]], bsize=13.5)
card(s, ML+half+0.3, 1.85, half, 2.6, "Output",
     [[("For every detected tree: a bounding box, an instance mask, a confidence score, crown area and lat/lng.", 13.5, INK2, False, False)],
      [("Tiles are stitched, duplicates merged, results shown on a map with per-area counts and ", 13.5, INK2, False, False),
       ("GeoJSON / CSV / HTML", 13.5, INK, True, False),(" export.", 13.5, INK2, False, False)]], bsize=13.5)
card(s, ML, 4.7, CW, 0.95, "",
     [[("capture  →  tile  →  model f(·)  →  { box, mask, conf, area, lat/lng }  →  merge & map & export", 14, INK, True, False)]], accent=BLUE, bsize=14)
footer(s, "Criterion 1 · Problem formulation", 7)

# 08 LITERATURE REVIEW (table, gap column, no our row)
s = slide(); band(s); title(s, "Literature review — what others report, and the gap each leaves", y=1.0, size=24)
lit_hdr = ["Method","Authors · year","Data","Best metric","Gap for our problem"]
lit = [
 ["Mask R-CNN (MCAN)","Lv et al. · 2023","UAV RGB, China","Det AP 92.4%","Forest UAV, not urban satellite; research result only"],
 ["RetinaNet / Faster R-CNN","dos Santos et al. · 2019","UAV RGB, Brazil","AP 82–93%","Detection only — no segmentation, no system"],
 ["YOLOv12m","Abbas & Damaševičius · 2025","RGB satellite","mAP@50 90.8%","Benchmark on a clean public set; no deployable tool"],
 ["DeepLabV3+ / U-Net","Martins 2021 / Wang 2021","Aerial 10–32 cm","F1 91% · IoU 96%","Semantic mask, not per-tree; not Central Asia"],
 ["DeepForest (off→fine)","Ventura et al. · 2024","NAIP 60 cm, USA","F 0.42 → 0.73","US imagery; collapses without local fine-tuning"],
 ["YOLOv5x","Velasquez-Camacho · 2023","Ground-level, Spain","F1 84.9%","Europe, ground-level; no GIS/coordinate pipeline"],
 ["DeepForest (urban)","Dakov & Petrova-Antonova · 2024","Aerial, Sofia","F1 0.67–0.69","Closest analogue — but Europe, still research-only"],
]
cols = [2.05, 2.35, 1.75, 1.55, 4.42]
rows = len(lit)+1
tbl_h = 3.5
gfx = s.shapes.add_table(rows, 5, Inches(ML), Inches(1.65), Inches(sum(cols)), Inches(tbl_h)).table
gfx.first_row = False; gfx.horz_banding = False
for j,wd in enumerate(cols): gfx.columns[j].width = Inches(wd)
def setcell(c, text, size, col, bold, fill, align=PP_ALIGN.LEFT, ital=False):
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.margin_left=Pt(6); c.margin_right=Pt(5); c.margin_top=Pt(3); c.margin_bottom=Pt(3)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; f=r.font; f.size=Pt(size); f.bold=bold; f.italic=ital; f.name=SERIF; f.color.rgb=col
for j,htxt in enumerate(lit_hdr):
    setcell(gfx.cell(0,j), htxt, 11, WHITE, True, NAVY)
for i,row in enumerate(lit, start=1):
    fill = WHITE if i%2 else RGBColor(0xF6,0xF9,0xFC)
    for j,val in enumerate(row):
        if j==0: setcell(gfx.cell(i,j), val, 10.5, INK, True, fill)
        elif j==3: setcell(gfx.cell(i,j), val, 10.5, BLUE, True, fill)
        elif j==4: setcell(gfx.cell(i,j), val, 10.5, INK2, False, fill, ital=True)
        else: setcell(gfx.cell(i,j), val, 10.5, INK2, False, fill)
gfx.rows[0].height = Inches(0.34)
rich(s, ML, 5.35, CW, 0.9, [
 [("Every study reports high accuracy — on ", 13, INK2, False, True),("its own region and sensor", 13, INK, True, False),
  (", as a research result. Overall limitation: ", 13, INK2, False, True),
  ("none was tested on Central Asia, and none delivers a deployable, operator-run system.", 13, INK, True, False),
  ("  That gap — not a higher mAP — is what this project fills. (31 references reviewed.)", 13, INK2, False, True)]], leading=1.15)
footer(s, "Criterion 1 · Critical literature analysis · ≥ 20 sources", 8)

# 09 ARCHITECTURE
s = slide(); band(s); eyebrow(s, "The core of the project"); title(s, "The automated pipeline", y=1.18)
blocks = [("Front end — React 18 + Leaflet", RGBColor(0x14,0x31,0x5A), False),
          ("API — FastAPI", RGBColor(0x1E,0x5E,0x32), False),
          ("Sliding-window tiling (640 + 128)", RGBColor(0x2E,0x7D,0x32), False),
          ("Engine — swappable (YOLO · M-RCNN · DF · Ensemble)", RGBColor(0x17,0xA2,0xC9), True),
          ("Global NMS — cross-tile merge", RGBColor(0x3F,0x86,0xC6), False),
          ("Geo — pixel → WGS-84", RGBColor(0x3F,0x6F,0xB5), False),
          ("SQLite store → Map & export", RGBColor(0x4F,0x86,0xC6), False)]
bx, bw, by0, bh, gap = ML, 4.7, 2.05, 0.5, 0.12
for i,(t,c,hl) in enumerate(blocks):
    block(s, bx, by0+i*(bh+gap), bw, bh, t, c, outline=hl)
rx = bx+bw+0.45; rw = SW-MR-rx
card(s, rx, 2.2, rw, 1.35, "One automated chain",
     [[("Capture → tile → detect → merge → georeference → store → map → export. ", 13.5, INK2, False, False),
       ("No manual step", 13.5, INK, True, False),(" between the operator drawing an area and the finished, downloadable inventory.", 13.5, INK2, False, False)]],
     accent=CYAN, bsize=13.5)
card(s, rx, 3.7, rw, 1.35, "Swappable engine — the key design idea",
     [[("The detection model is just one plug-in step. Swap YOLO for Mask R-CNN, DeepForest or an ensemble without touching the rest — a better model tomorrow drops straight in.", 13.5, INK2, False, False)]],
     accent=BLUE, bsize=13.5)
card(s, rx, 5.2, rw, 0.95, "Performance",
     [[("~1 km² at zoom 19 in ~18 s on a single laptop GPU, with live streaming progress.", 13.5, INK2, False, False)]],
     accent=CYAN, bsize=13.5)
footer(s, "Criterion 1–2 · System architecture + software", 9)

# 10 BEYOND THE MINIMUM
s = slide(); band(s); eyebrow(s, "Engineering effort")
title(s, "A minimal version would do far less — we built a tool", y=1.18, size=27)
vh = 3.15; vy = 2.15; vw = (CW-0.0)/2
# min side
rect(s, ML, vy, vw, vh, GREY, line=CARDLN)
tb(s, ML+0.22, vy+0.16, vw-0.4, 0.3, [("THE BARE MINIMUM", 13, MUTED, True, False)])
minl = ["–  One model, hard-coded","–  Just a segmentation mask","–  No coordinates, no export","–  No metrics, no comparison","–  Run once, nothing saved"]
rich(s, ML+0.22, vy+0.62, vw-0.44, vh-0.74, [[(x,14,INK2,False,False)] for x in minl], leading=1.3)
# built side
rect(s, ML+vw, vy, vw, vh, BUILT, line=CARDLN)
tb(s, ML+vw+0.22, vy+0.16, vw-0.4, 0.3, [("WHAT CANOPY ACTUALLY DOES", 13, NAVY, True, False)])
builtl = [("+  Two tile providers"," — ESRI & Google, fetched in-browser"),
          ("+  Model choice"," — 7 YOLO sizes · Mask R-CNN · DeepForest ± SAM 2"),
          ("+  Ensembles"," — Weighted Box Fusion + cross-YOLO vote"),
          ("+  Manage scans"," — name · hide · delete (cascade)"),
          ("+  Live metrics"," — tree count · canopy coverage % · crown area"),
          ("+  Geo + export"," — WGS-84 · GeoJSON / CSV / HTML · persistent map")]
rich(s, ML+vw+0.22, vy+0.62, vw-0.44, vh-0.74, [[(a,12.5,NAVY,True,False),(b,12.5,NAVY,False,False)] for a,b in builtl], leading=1.16)
rich(s, ML, vy+vh+0.22, CW, 0.6, [
 [("None of these is required to “detect a tree” — each is a deliberate choice to make the system ", 14, INK2, False, True),
  ("usable, comparable and repeatable.", 14, INK, True, False),
  (" That is where most of the engineering effort went.", 14, INK2, False, True)]], leading=1.1)
footer(s, "Criterion 2 · Depth & completeness of the software", 10)

# 11 DATA / DB / HARDWARE
s = slide(); band(s); eyebrow(s, "A component of the system — inputs & infrastructure")
title(s, "Data, database & hardware", y=1.18)
lw = 6.2
card(s, ML, 2.15, lw, 1.5, "Dataset",
     [[("High-zoom Astana captures (~0.3 m/px). We hand-labelled ", 13, INK2, False, False),
       ("~5,500 tree crowns", 13, INK, True, False),(" across ", 13, INK2, False, False),("~100 images", 13, INK, True, False),
       (" as instance polygons (~8,700 tiles after 640+128 tiling). Held-out common set: ", 13, INK2, False, False),
       ("14 images · 702 trees", 13, INK, True, False),(".", 13, INK2, False, False)]], bsize=13)
card(s, ML, 3.8, lw, 1.7, "Database — SQLite",
     [[("scan_sessions → snapshots → runs → detections", 13.5, INK, True, False)],
      [("One-to-many, ON DELETE CASCADE — delete a scan, its detections go too.", 13, INK2, False, False)]], bsize=13)
rx = ML+lw+0.4; rw = SW-MR-rx
card(s, rx, 2.15, rw, 1.45, "Hardware",
     [[("RTX 4060 8 GB — YOLO + all common-set inference · RTX 4070 — Mask R-CNN · RTX 4050 — DeepForest. The 8 GB ceiling forces batch 2 + mixed precision — deliberately modest, to prove it runs ", 12.5, INK2, False, False),
       ("without a cluster.", 12.5, INK, True, False)]], bsize=12.5)
_,iw,ih = place_img(s, "yolo_v2_finetune_val_4tile_strip.png", rx, 3.75, rw)
footer(s, "Criterion 1–2 · DB schema + hardware", 11)

# 12 MATHS
s = slide(); band(s); eyebrow(s, "The maths, in plain terms")
title(s, "Three simple ideas — one comparison number", y=1.18, size=28)
half=(CW-0.3)/2
card(s, ML, 2.15, half, 1.55, "1 · IoU — did the box land on the tree?",
     [[("IoU = (overlap area) ÷ (combined area). 1.0 = perfect, 0 = miss. A tree counts as found when ", 13, INK2, False, False),("IoU ≥ 0.5.", 13, INK, True, False)]], bsize=13)
card(s, ML, 3.85, half, 1.55, "2 · Precision & recall",
     [[("Precision — of the trees we flagged, how many were real. Recall — of all real trees, how many we found: ", 13, INK2, False, False),("ours ≈ 30%", 13, BLUE, True, False),(" at the default setting.", 13, INK2, False, False)]], bsize=13)
card(s, ML+half+0.3, 2.15, half, 1.55, "3 · mAP@50 — the single score",
     [[("Sweep the confidence threshold, track precision against recall, take the area under that curve (IoU ≥ 0.5). One honest number to compare every engine — the ", 13, INK2, False, False),("0.315", 13, INK, True, False),(" on the results slide.", 13, INK2, False, False)]], bsize=13)
card(s, ML+half+0.3, 3.85, half, 1.55, "How the two detectors work",
     [[("YOLOv8 — one pass: grid → boxes + masks in a single look (fast). Mask R-CNN — two passes: propose regions, then refine (slower, careful).", 13, INK2, False, False)]], bsize=13)
footer(s, "Criterion 1–2 · Mathematical support (≤ 10 pts)", 12)

# 13 ENGINE DIVIDER
s = slide(); band(s); eyebrow(s, "Four interchangeable engines · one pipeline", y=2.2)
tb(s, ML, 2.6, CW-2.0, 1.8, [("The system doesn't depend on one model — any of four can drive it.", 38, BLUE, True, False)], leading=1.05)
tb(s, ML, 4.7, CW, 0.8, [("YOLOv8-seg (Rasul) · Mask R-CNN (Berik) · DeepForest + SAM 2 (Anuar) · plus two ensembles. Each plugs into the same automated pipeline and is scored on the same 14 images.", 16, INK2, False, False)])
footer(s, "Team diploma — each member presents & is scored on their branch", 13)

# 14 YOLO
s = slide(); band(s); eyebrow(s, "Engine #1 · Rasul Aidarkhanov · 23 experiments")
title(s, "YOLOv8x-seg — from 0.131 to 0.315", y=1.18, size=28)
_,iw,ih = place_img(s, "yolo_all_runs_map50.png", ML, 2.2, 6.7)
tb(s, ML, 2.2+ih+0.06, 6.7, 0.3, [("Box mAP@50 across every experiment on the common set.", 12, MUTED, False, True)])
rx = ML+7.0; rw = SW-MR-rx
tb(s, rx, 2.1, rw, 0.9, [("+140%", 64, BLUE, True, False)])
tb(s, rx, 3.05, rw, 0.4, [("Box mAP@50 over my own first model (0.131 → 0.315).", 13, INK2, False, False)])
card(s, rx, 3.55, rw, 1.05, "What actually mattered",
     [[("① Start from pretrained (COCO) weights · ② tile at the resolution you deploy at · ③ keep augmentation moderate.", 12.5, INK2, False, False)]], bsize=12.5)
card(s, rx, 4.75, rw, 1.1, "What I also found",
     [[("Default augmentation beat hand-tuned; model size is U-shaped; run-to-run noise ≈ ±0.03 — so small gaps aren't real.", 12.5, INK2, False, False)]], bsize=12.5)
footer(s, "Criterion 1 · Results + ablation · YOLOv8x-seg v4 · M14", 14)

# 15 MASK R-CNN
s = slide(); band(s); eyebrow(s, "Engine #2 · Berik Sharipov")
title(s, "Mask R-CNN — the careful, mask-first engine", y=1.18, size=27)
tb(s, ML, 2.0, CW, 0.7, [("A classic two-stage, region-based detector. Where YOLO looks once, Mask R-CNN looks twice — propose, then refine — slower, but with naturally crisp instance masks. One of the swappable engines behind the same pipeline.", 14.5, INK2, False, False)], leading=1.2)
half=(CW-0.3)/2
card(s, ML, 2.95, half, 1.4, "How it works", [[("An RPN suggests candidate boxes → RoIAlign crops each region → three heads output a box, a class and a per-pixel mask. Segmentation is built in, so outlines hug the crown.", 12.5, INK2, False, False)]], bsize=12.5)
card(s, ML, 4.5, half, 1.4, "How Berik trained it", [[("Transfer-learned from a COCO-pretrained ResNet-50-FPN, fine-tuned on the same Astana tiles (v2 + v3), single “tree” class, same split. 50 epochs, SGD.", 12.5, INK2, False, False)]], bsize=12.5)
card(s, ML+half+0.3, 2.95, half, 1.4, "How it does on Astana", [[("Box mAP@50 0.166, Mask 0.158. At conf 0.5: precision 0.44 · recall 0.22. Strongest on medium, well-separated crowns.", 12.5, INK2, False, False)]], bsize=12.5)
card(s, ML+half+0.3, 4.5, half, 1.4, "Why it stays in the app", [[("Below YOLO overall — but it recovers some larger trees the others miss, and its masks are the cleanest. So the pipeline keeps it as a selectable engine.", 12.5, INK2, False, False)]], accent=CYAN, bsize=12.5)
footer(s, "Team member · Berik Sharipov · box 0.166 · mask 0.158", 15)

# 16 DEEPFOREST + SAM2
s = slide(); band(s); eyebrow(s, "Engine #3 · Anuar Totin")
title(s, "DeepForest + SAM 2 — the forest specialist, re-taught", y=1.18, size=26)
tb(s, ML, 2.0, CW, 0.7, [("A two-part engine: DeepForest (a RetinaNet-style crown detector) finds the trees, then SAM 2 turns each box into a high-quality mask. It produced the 0.012 that started the whole project.", 14.5, INK2, False, False)], leading=1.2)
card(s, ML, 2.95, half, 1.4, "How it works", [[("DeepForest outputs a box per crown; each box becomes a prompt for SAM 2, which segments the exact pixels — zero-shot, no extra training.", 12.5, INK2, False, False)]], bsize=12.5)
card(s, ML, 4.5, half, 1.4, "The domain problem", [[("DeepForest is pretrained on North-American forest imagery. Pointed at a dry steppe city it was nearly blind — 0.012. A mismatch, not a bug.", 12.5, INK2, False, False)]], bsize=12.5)
card(s, ML+half+0.3, 2.95, half, 1.4, "How Anuar fixed it", [[("Fine-tuned DeepForest on the Astana data (LR 1e-4, 30 ep) + SAM 2 masks. Result: 0.012 → 0.146 Box mAP@50 (Mask 0.134) — a ×12 jump.", 12.5, INK2, False, False)]], bsize=12.5)
card(s, ML+half+0.3, 4.5, half, 1.4, "Why it stays in the app", [[("SAM 2 gives the best general-purpose masks of any engine, and the off-the-shelf → fine-tuned story is the project's core lesson.", 12.5, INK2, False, False)]], accent=CYAN, bsize=12.5)
footer(s, "Team member · Anuar Totin · 0.012 → 0.146 · ×12", 16)

# 17 WHY SEVERAL MODELS
s = slide(); band(s); eyebrow(s, "A system choice, not a model choice")
title(s, "Same garden, eight engines, eight answers", y=1.18, size=27)
_,iw,ih = place_img(s, "cross_model_8way_botanical.jpg", (SW-7.2)/2, 2.05, 7.2)
yb = 2.05+ih+0.12
half=(CW-0.4)/2
tb(s, ML, yb, half, 0.9, [("The same Botanical Garden, scanned by eight engines — each outlines a different set of trees and reports a different count (672–819). None is simply “right”.", 13.5, INK2, False, False)], leading=1.18)
tb(s, ML+half+0.4, yb, half, 0.9, [("So the system lets the operator switch engines and combine them — Weighted Box Fusion + a cross-YOLO 4-vote that drops single-model false positives. Choice is a feature.", 13.5, INK2, False, False)], leading=1.18)
footer(s, "Criterion 1 · Own contribution — design rationale · WBF · cross-YOLO vote", 17)

# 18 RESULTS
s = slide(); band(s); eyebrow(s, "Head-to-head · same 14 images · 702 trees")
title(s, "One common test set, every model", y=1.18, size=28)
res = [["YOLOv8x-seg v4 (ours)","0.315","0.289"],["YOLOv8x-seg v3","0.268","0.244"],
       ["Mask R-CNN","0.166","0.158"],["DeepForest + SAM 2","0.146","0.134"],
       ["YOLO v1 baseline","0.131","0.134"],["NEON DeepForest (off-the-shelf)","0.012","—"]]
rcols=[3.3,1.4,1.4]; rt = s.shapes.add_table(len(res)+1,3, Inches(ML), Inches(2.15), Inches(sum(rcols)), Inches(2.7)).table
rt.first_row=False; rt.horz_banding=False
for j,wd in enumerate(rcols): rt.columns[j].width=Inches(wd)
for j,h in enumerate(["Model","Box mAP@50","Mask mAP@50"]): setcell(rt.cell(0,j), h, 11, WHITE, True, NAVY)
for i,row in enumerate(res, start=1):
    ours = (i==1); fill = BUILT if ours else (WHITE if i%2 else RGBColor(0xF6,0xF9,0xFC))
    for j,v in enumerate(row):
        col = NAVY if ours else (INK if j==0 else BLUE)
        setcell(rt.cell(i,j), v, 11, col, (j==0 or ours), fill)
rich(s, ML, 5.05, 6.3, 1.0, [
 [("Champion operating point (conf 0.25): ", 12.5, INK2, False, True),("precision 0.52 · recall 0.29.", 12.5, INK, True, False),
  (" Margins between top models are small — partly within run-to-run noise.", 12.5, INK2, False, True)]], leading=1.15)
_,iw,ih = place_img(s, "model_comparison_barchart.png", ML+6.6, 2.15, SW-MR-(ML+6.6))
tb(s, ML+6.6, 2.15+ih+0.08, SW-MR-(ML+6.6), 0.8, [("These engines are in the same ballpark on Astana — what matters is that any of them can drive the system.", 12, MUTED, False, True)], leading=1.15)
footer(s, "Criterion 1–2 · Results + comparison · M14 · 702 polygons", 18)

# 19 CANOPY OVERVIEW
s = slide(); band(s); eyebrow(s, "The deliverable")
title(s, "Canopy — the system in use", y=1.18)
_,iw,ih = place_img(s, "ui_canopy_map_view.png", ML, 2.15, 6.9)
rx = ML+7.2; rw = SW-MR-rx
card(s, rx, 2.15, rw, 1.5, "One flow", [[("Select an area of Astana → the app tiles it, runs the chosen model, and draws every detected tree on the map — confidence filter, Point / Box / Polygon / Heat layers.", 13, INK2, False, False)]], bsize=13)
card(s, rx, 3.8, rw, 1.5, "Grows as you use it", [[("Every scan is stored; the city map is the running aggregate of everything ever processed. Coverage and counts update live.", 13, INK2, False, False)]], bsize=13)
tb(s, rx, 5.4, rw, 0.5, [("Counts in screenshots come from repeated demo scans — not a full city census.", 11.5, MUTED, False, True)], leading=1.1)
footer(s, "Criterion 1–2 · Information system built & working", 19)

# 20 CANOPY FEATURES
s = slide(); band(s); eyebrow(s, "The choices that make it a tool, not a demo")
title(s, "Choose · combine · manage · measure", y=1.18, size=28)
iw3 = (CW-2*0.25)/3
imgs3 = ["ui_canopy_scan_model_modal.png","ui_canopy_library_modal.png","ui_canopy_map_heat.png"]
ih3=0
for i,f in enumerate(imgs3):
    _,_,ih3 = place_img(s, f, ML+i*(iw3+0.25), 2.1, iw3)
cy = 2.1+ih3+0.18
feats=[("Choose & combine","7 YOLO sizes · Mask R-CNN · DeepForest ± SAM 2 · WBF & cross-YOLO vote — chosen per scan, with provider (ESRI / Google) and confidence."),
       ("Manage","Every scan & snapshot is renamable, hideable and deletable (cascade) from a library modal — the workspace stays usable over time."),
       ("Measure","Density heat-map, per-area count, green-space coverage % and crown area inside a drawn polygon — metrics, not just dots.")]
for i,(t,b) in enumerate(feats):
    card(s, ML+i*(iw3+0.25), cy, iw3, 1.55, t, [[(b,12,INK2,False,False)]], tsize=13.5, bsize=12)
footer(s, "Criterion 2 · Feature depth of the system", 20)

# 21 DEMO
s = slide(); band(s); eyebrow(s, "The pipeline running end-to-end")
title(s, "Live demo — automated scan of the Botanical Garden", y=1.18, size=26)
poster = os.path.join(ASSET, "canopy_demo_poster.jpg")
pic = s.shapes.add_picture(poster, Inches(ML), Inches(2.2), width=Inches(6.9))
pic.line.color.rgb=LINE; pic.line.width=Pt(1)
# play glyph
pl = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(ML+3.1), Inches(3.35), Inches(0.7), Inches(0.7))
pl.rotation=90; pl.fill.solid(); pl.fill.fore_color.rgb=WHITE; _noline(pl); pl.shadow.inherit=False
rx = ML+7.2; rw = SW-MR-rx
card(s, rx, 2.2, rw, 1.2, "What you're seeing", [[("The v4 engine scanning the Botanical Garden end-to-end in Canopy: tile → detect → map, every tree drawn live.", 12.5, INK2, False, False)]], bsize=12.5)
card(s, rx, 3.55, rw, 1.1, "Why this scene", [[("A dense, well-defined planting — the clearest case. Honest reminder: on sparser scenes recall drops.", 12.5, INK2, False, False)]], bsize=12.5)
card(s, rx, 4.8, rw, 1.05, "Full recording", [[("drive.google.com/file/d/1OceaLC4yMEgLtIThVR8t7_B4VRRPNaNf/view", 11, BLUE, False, False)]], accent=CYAN, bsize=11)
tb(s, ML, 2.2+6.9/1.76+0.05, 6.9, 0.3, [("Replace this poster with the embedded video in Canva (assets/canopy_demo.mp4).", 10.5, MUTED, False, True)])
footer(s, "Criterion 2 · Demonstration of the working system", 21)

# 22 CONCLUSION
s = slide(); band(s); title(s, "Conclusion", y=1.0)
conc=[("O1 · Dataset","Custom Astana dataset built from scratch — ~100 images, ~5,500 crown polygons (~8,700 tiles); held-out M14 = 14 images / 702 trees.", BLUE, CARD),
      ("O2 · Detection engines","Four engines trained & compared; YOLOv8x-seg leads at Box mAP@50 = 0.315 (+140% over our v1 baseline).", BLUE, CARD),
      ("O3 · Ensembles","WBF + cross-YOLO vote implemented in the system — choice and combination exposed to the operator.", BLUE, CARD),
      ("O4 · Fair evaluation","Every engine scored on one common set — the first such measurement for Astana (off-the-shelf ≈ 0.012).", BLUE, CARD),
      ("O5 · The system — Canopy","FastAPI + React + SQLite + Leaflet, deployed: provider & model choice, ensembles, manage, live metrics, geo + export.", CYAN, CARD),
      ("The headline","Not a score — a manual task is now an automated system. Accuracy is modest and honestly so; the system is the contribution.", BLUE, BUILT)]
cw2=(CW-0.4)/2; rh=1.3
for i,(t,b,ac,fl) in enumerate(conc):
    cx = ML + (i%2)*(cw2+0.4); cy = 1.8 + (i//2)*(rh+0.2)
    card(s, cx, cy, cw2, rh, t, [[(b,12.5,INK2,False,False)]], accent=ac, fill=fl, tsize=14, bsize=12.5)
footer(s, "Criterion 1 · Conclusions map 1:1 to tasks · manual → automated system", 22)

# 23 THANK YOU
s = slide()
# right-half image, white left
img = os.path.join(FIG, "ui_canopy_botanical_park.jpg")
p = s.shapes.add_picture(img, Inches(SW*0.5), 0, height=prs.slide_height)
band(s)
tb(s, ML, 2.0, SW*0.5-0.4, 2.2, [("Thank you — we'll gladly answer your questions.", 40, BLUE, True, False)], leading=1.05)
tb(s, ML, 4.4, SW*0.5-0.4, 1.2, [("From a manual, one-by-one task to an automated tree-mapping system for Astana — a working tool anyone can run.", 17, INK2, False, False)], leading=1.25)
tb(s, ML, 5.7, SW*0.5-0.4, 0.5, [("Rasul Aidarkhanov · Berik Sharipov · Anuar Totin   |   Supervisor: S. Satbayev", 13, MUTED, False, False)])

out = os.path.join(HERE, "deck_aitu_system.pptx")
prs.save(out)
print("Saved", out, "—", len(prs.slides._sldIdLst), "slides")
