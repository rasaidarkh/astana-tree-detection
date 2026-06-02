# -*- coding: utf-8 -*-
"""Diagram pack (AITU style, Raikhan-like) for the Canopy system deck.
Four editable architecture diagrams:
  1. System architecture — horizontal swim-lanes
  2. Tiled-inference pipeline — horizontal flow + config
  3. Dataset pipeline — medallion-style stages
  4. Four detection engines — numbered columns + adapter bar
Render: venv/Scripts/python.exe thesis/slides/build_diagrams.py
"""
import os, sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))

BLUE=RGBColor(0x2E,0x6D,0xB4); CYAN=RGBColor(0x29,0xAB,0xE2); NAVY=RGBColor(0x17,0x3E,0x63)
INK=RGBColor(0x1A,0x1A,0x1A); INK2=RGBColor(0x33,0x41,0x4D); MUTED=RGBColor(0x6A,0x76,0x82)
CARD=RGBColor(0xEE,0xF4,0xFB); CARDLN=RGBColor(0xCF,0xE0,0xF1); LINE=RGBColor(0xD7,0xDE,0xE4)
WHITE=RGBColor(0xFF,0xFF,0xFF); BUILT=RGBColor(0xDC,0xEB,0xFA)
GREEN=RGBColor(0x2E,0x7D,0x32); TEAL=RGBColor(0x16,0x97,0xA6)
PURPLE=RGBColor(0x6A,0x3F,0xA6)
ORANGE=RGBColor(0xE8,0x77,0x22); AMBER=RGBColor(0xF0,0x9A,0x16); GREY=RGBColor(0x5C,0x6B,0x7A)
SERIF="Times New Roman"
SW,SH=13.333,7.5; ML,MR=0.5,0.5; CW=SW-ML-MR

prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def _nl(sp): sp.line.fill.background()

def rect(s,x,y,w,h,fill,line=None,lw=0.75,rounded=False):
    shape=MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.shadow.inherit=False
    if line is None: _nl(sp)
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    return sp

def txt_in(sp,text,size,col,bold=False,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE):
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(4);tf.margin_right=Pt(4);tf.margin_top=Pt(2);tf.margin_bottom=Pt(2)
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; f=r.font; f.size=Pt(size); f.bold=bold; f.name=SERIF; f.color.rgb=col
    return sp

def tb(s,x,y,w,h,paras,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,leading=1.05):
    box=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,runs in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=leading; p.space_after=Pt(3); p.space_before=Pt(0)
        for (t,sz,c,b,it) in runs:
            r=p.add_run(); r.text=t; f=r.font; f.size=Pt(sz); f.bold=b; f.italic=it; f.name=SERIF; f.color.rgb=c
    return box

def band(s):
    b=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(0.6)); b.shadow.inherit=False; _nl(b)
    try:
        b.fill.gradient(); st=b.fill.gradient_stops
        st[0].color.rgb=NAVY; st[0].position=0.0; st[1].color.rgb=CYAN; st[1].position=1.0
        b.fill.gradient_angle=0.0
    except Exception:
        b.fill.solid(); b.fill.fore_color.rgb=NAVY
    rect(s,SW-1.55,0,1.55,0.6,CYAN)
    rect(s,0.62,0,0.8,0.78,CYAN)
    tri=s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,Inches(0.62),Inches(0.78),Inches(0.8),Inches(0.30))
    tri.rotation=180; tri.fill.solid(); tri.fill.fore_color.rgb=CYAN; tri.shadow.inherit=False; _nl(tri)
    ring=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.82),Inches(0.12),Inches(0.4),Inches(0.4))
    ring.fill.background(); ring.line.color.rgb=WHITE; ring.line.width=Pt(2); ring.shadow.inherit=False
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.96),Inches(0.30),Inches(0.12),Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb=WHITE; _nl(dot); dot.shadow.inherit=False

def eyebrow(s,t,y=0.86): tb(s,ML,y,CW,0.3,[[(t,13,BLUE,True,False)]])
def title(s,t,y=1.12,size=29): tb(s,ML,y,CW,0.7,[[(t,size,INK,True,False)]])
def footer(s,left,page):
    rect(s,ML,SH-0.6,CW,0.012,LINE)
    tb(s,ML,SH-0.53,CW-0.6,0.3,[[(left,9,MUTED,False,True)]])
    tb(s,SW-MR-0.6,SH-0.53,0.6,0.3,[[(str(page),9,MUTED,False,True)]],align=PP_ALIGN.RIGHT)

def arrow(s,x,y,col=CYAN,w=0.32,h=0.34):
    a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x),Inches(y),Inches(w),Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb=col; _nl(a); a.shadow.inherit=False
    return a

def numcircle(s,x,y,n,col,d=0.34):
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb=col; _nl(c); c.shadow.inherit=False
    txt_in(c,str(n),14,WHITE,bold=True)
    return c

# ============================================================== 1. SYSTEM ARCHITECTURE
s=slide(); band(s); eyebrow(s,"The system, end to end")
title(s,"System architecture — from satellite tile to GIS export",size=26)
# caption strip
rect(s,ML,1.7,CW,0.34,NAVY)
txt_in(rect(s,ML,1.7,CW,0.34,NAVY),"Canopy · automated tree-inventory pipeline — no manual step between drawing an area and the downloadable inventory",11,WHITE,bold=True)

lanes=[
 ("INPUT",NAVY,["ESRI World Imagery tiles","Google Satellite tiles","File upload — PNG / JPG / GeoTIFF","Or: draw an area on the map"]),
 ("FRONTEND — React 18 + Leaflet",BLUE,["Pick area · model · confidence","Auto-Zoom / Polygon scan","Live NDJSON progress","Point / Box / Polygon / Heat"]),
 ("API — FastAPI",GREEN,["capture_bbox()","sliding-window tile (640+128)","batch inference","Global NMS cross-tile merge"]),
 ("MODEL LAYER — swappable",CYAN,["YOLOv8-seg (7 sizes)","Mask R-CNN","DeepForest + SAM 2","Ensembles: WBF · cross-YOLO"]),
 ("GEO + STORAGE",BLUE,["pixel → WGS-84 (4 modes)","SQLite persistence","sessions→snapshots→runs→detections","ON DELETE CASCADE"]),
 ("OUTPUT",NAVY,["City map + density heat-map","Per-area count · coverage %","GeoJSON · CSV · standalone HTML","(QGIS / ArcGIS ready)"]),
]
n=len(lanes); gap=0.16; aw=0.30
lw=(CW-(n-1)*gap-(n-1)*aw)/n
ly=2.25; lh=3.95; hh=0.62
x=ML
hi_idx=3
for i,(hd,col,items) in enumerate(lanes):
    # body
    body=rect(s,x,ly,lw,lh,WHITE,line=(NAVY if i==hi_idx else CARDLN),lw=(1.6 if i==hi_idx else 0.8))
    # header
    rect(s,x,ly,lw,hh,col)
    txt_in(rect(s,x,ly,lw,hh,col),hd,11.5,WHITE,bold=True)
    # items as small boxes
    iy=ly+hh+0.12; ih=(lh-hh-0.24-(len(items)-1)*0.1)/len(items)
    for it in items:
        bx=rect(s,x+0.1,iy,lw-0.2,ih,CARD,line=CARDLN,lw=0.5)
        txt_in(bx,it,9.7,INK2,align=PP_ALIGN.CENTER)
        iy+=ih+0.1
    x+=lw
    if i<n-1:
        arrow(s,x+(gap+aw)/2-aw/2,ly+lh/2-0.17,CYAN,w=aw,h=0.34); x+=gap+aw
tb(s,ML,6.4,CW,0.4,[[("Six layers, one automated chain. The detection model is just ",12,INK2,False,False),
    ("one swappable layer",12,BLUE,True,False),(" — a better model drops in without touching capture, tiling, geo or storage.",12,INK2,False,False)]],leading=1.1)
footer(s,"Criterion 1–2 · System architecture + software · FastAPI · React · SQLite · Leaflet",9)

# ============================================================== 2. TILED INFERENCE PIPELINE
s=slide(); band(s); eyebrow(s,"How one scan is processed")
title(s,"Tiled-inference pipeline",size=29)
steps=["Capture\n(~1 km², z19)","Sliding-window\ntile · 640 + 128","Per-tile\nYOLOv8-seg","Translate to\nglobal pixels","Global NMS\ndedup · IoU 0.5","Instance polygons\n+ confidence"]
cols=[NAVY,GREEN,CYAN,BLUE,RGBColor(0x3F,0x6F,0xB5),RGBColor(0x4F,0x86,0xC6)]
n=len(steps); gap=0.12; aw=0.26
bw=(CW-(n-1)*gap-(n-1)*aw)/n
by=2.35; bh=1.3
x=ML
for i,(st,col) in enumerate(zip(steps,cols)):
    rect(s,x,by,bw,0.1,col)  # top accent
    bx=rect(s,x,by,bw,bh,WHITE,line=CARDLN,lw=0.9)
    rect(s,x,by,bw,0.1,col)
    txt_in(bx,st.replace("\n","\n"),12,INK,bold=True)
    x+=bw
    if i<n-1:
        arrow(s,x+gap/2,by+bh/2-0.16,CYAN,w=aw,h=0.32); x+=gap+aw
# config callouts
cy=4.2; ch=1.7; half=(CW-0.3)/2
def card(x,y,w,h,ttl,body,accent=BLUE):
    rect(s,x,y,w,h,CARD,line=CARDLN); rect(s,x,y,0.06,h,accent)
    tb(s,x+0.18,y+0.13,w-0.34,h-0.26,[[(ttl,14,BLUE,True,False)],[(body,13,INK2,False,False)]],leading=1.15)
card(ML,cy,half,ch,"Why tiling",
     "Zoom-19 captures are far larger than 640 px. A sliding window with 128-px overlap lets the model see every crown at native resolution — small street trees aren't shrunk away.")
card(ML+half+0.3,cy,half,ch,"Why Global NMS",
     "A crown on a tile boundary is detected twice. After translating all boxes to global coordinates, a greedy NMS pass (IoU > 0.5) keeps the most confident copy — one tree, one detection.")
tb(s,ML,6.15,CW,0.4,[[("Same pattern reused at train time (tile the dataset) and at inference time (tile the scan) — identical 640 + 128 geometry.",12,MUTED,False,True)]])
footer(s,"Criterion 1–2 · Inference pipeline · 640 + 128 sliding window",10)

# ============================================================== 3. DATASET PIPELINE (medallion)
s=slide(); band(s); eyebrow(s,"A component of the system — how the data was built")
title(s,"Dataset pipeline — built from scratch",size=28)
stages=[
 ("CAPTURE",ORANGE,["Google Earth Pro + ESRI","screenshots, zoom 17–19","≈ 0.3 m / pixel","~100 source images","mixed districts of Astana"]),
 ("ANNOTATE",GREY,["CVAT polygon labels","single class — “tree”","model-in-the-loop pre-label","~5 500 crown polygons","≈ 25 → 4 min / image"]),
 ("TILE & SPLIT",AMBER,["640 + 128 sliding window","≈ 8 700 tiled instances","COCO → YOLO format","source-level train / val split","held-out M14 = 14 img / 702 trees"]),
]
n=len(stages); aw=0.5
cw=(CW-(n-1)*aw-(n-1)*0.2)/n
cy=2.3; chh=3.4; hh=0.6
x=ML
for i,(hd,col,items) in enumerate(stages):
    rect(s,x,cy,cw,chh,WHITE,line=CARDLN,lw=0.9)
    rect(s,x,cy,cw,hh,col); txt_in(rect(s,x,cy,cw,hh,col),hd,15,WHITE,bold=True)
    iy=cy+hh+0.16
    tb(s,x+0.2,iy,cw-0.4,chh-hh-0.3,[[("▸  "+it,13,INK2,False,False)] for it in items],leading=1.35)
    x+=cw
    if i<n-1: arrow(s,x+0.1,cy+chh/2-0.18,BLUE,w=0.36,h=0.38); x+=aw+0.2
rect(s,ML,6.0,CW,0.55,BUILT,line=CARDLN); rect(s,ML,6.0,0.06,0.55,BLUE)
tb(s,ML+0.18,6.1,CW-0.36,0.4,[[("No public Astana tree dataset existed before this work — every polygon is hand-labelled. The dataset is reusable for future research in the region.",13,NAVY,True,False)]])
footer(s,"Criterion 1–2 · Dataset engineering · CVAT · ~100 imgs · ~5 500 crowns",11)

# ============================================================== 4. FOUR ENGINES (numbered columns)
s=slide(); band(s); eyebrow(s,"Four interchangeable engines · one adapter interface")
title(s,"The detection engines behind the pipeline",size=27)
eng=[
 (1,"YOLOv8-seg",BLUE,["One-stage, anchor-free","tile → boxes + masks in one pass","23-experiment ablation","Box mAP@50 = 0.315 (champion)"]),
 (2,"Mask R-CNN",TEAL,["Two-stage, region-based","RPN → RoIAlign → mask head","cleanest instance masks","Box 0.166 · Mask 0.158"]),
 (3,"DeepForest + SAM 2",PURPLE,["Crown detector + SAM 2 masks","off-the-shelf 0.012 on Astana","fine-tuned + zero-shot masks","0.012 → 0.146  (×12)"]),
 (4,"Ensembles",GREEN,["Weighted Box Fusion (YOLO+DF)","cross-YOLO 4-vote","drops single-model false positives","choice + combination in-app"]),
]
n=len(eng); cw=(CW-(n-1)*0.22)/n
cy=2.3; chh=3.5; hh=0.66
x=ML
for (num,nm,col,items) in eng:
    rect(s,x,cy,cw,chh,WHITE,line=CARDLN,lw=0.9)
    rect(s,x,cy,cw,hh,col)
    numcircle(s,x+0.14,cy+hh/2-0.17,num,WHITE if False else NAVY)  # dark circle on color header
    txt_in(rect(s,x+0.5,cy,cw-0.5,hh,col),nm,13.5,WHITE,bold=True,align=PP_ALIGN.LEFT)
    iy=cy+hh+0.16
    tb(s,x+0.18,iy,cw-0.34,chh-hh-0.3,[[("·  "+it,12,INK2,False,False)] for it in items],leading=1.4)
    x+=cw+0.22
rect(s,ML,6.05,CW,0.5,NAVY)
txt_in(rect(s,ML,6.05,CW,0.5,NAVY),"All four implement the same adapter interface — selected per scan from one model picker, scored on the same M14 set.",12.5,WHITE,bold=True)
footer(s,"Criterion 1–2 · Pluggable adapter interface · scored on M14",15)

out=os.path.join(HERE,"deck_diagrams.pptx")
prs.save(out)
print("Saved",out,"—",len(prs.slides._sldIdLst),"slides")
